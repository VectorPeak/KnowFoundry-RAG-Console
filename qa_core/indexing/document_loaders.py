"""LangChain 文档加载器注册表。新增文件类型只需维护注册项，无需 if/elif 分支。"""

from __future__ import annotations

from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol

import fitz
from langchain_core.documents import Document
from langchain_community.document_loaders import TextLoader
from docx import Document as DocxDocument
from pptx import Presentation

from qa_core.config.logging_config import get_logger
from qa_core.indexing.table_documents import load_table_file
logger = get_logger(__name__)

class SupportsLoad(Protocol):
    """文档加载器最小协议。不直接导入 LangChain BaseLoader，避免提前拉起 heavy 依赖。"""
    def load(self) -> list[Document]:
        """加载文件并返回 LangChain Document 列表。"""
        ...

LoaderFactory = Callable[[Path], SupportsLoad]

@dataclass(frozen=True)
class DocumentLoaderSpec:
    """文档加载器注册项。保存 factory 而非裸 loader class，因不同 loader 构造参数不同。"""

    suffixes: tuple[str, ...]
    factory: LoaderFactory
    description: str

    def create_loader(self, path: Path) -> SupportsLoad:
        """创建实际 loader 实例。

        工厂函数统一接收 Path，内部自行决定是否需要 encoding、mode 等参数。这样
        `load_file()` 不需要知道 TextLoader、PyPDFLoader、Docx2txtLoader 的构造差异。
        """
        return self.factory(path)

def _utf8_text_loader(path: Path) -> SupportsLoad:
    """创建 UTF-8 文本加载器。.md 也走 TextLoader 以保留 Markdown 标题给后续结构化切分。"""
    return TextLoader(str(path), encoding="utf-8")

def _pdf_loader(path: Path) -> SupportsLoad:
    """创建 PDF 文本层加载器。使用 PyMuPDF，适合有中文文本层的业务 PDF。"""
    return PyMuPdfLoader(path)

def _word_loader(path: Path) -> SupportsLoad:
    """创建 Word 文档加载器。优先使用当前依赖里的 python-docx，避免额外依赖 docx2txt。"""
    return PythonDocxLoader(path)


def _powerpoint_loader(path: Path) -> SupportsLoad:
    """创建 PPT/PPTX 文档加载器。优先使用当前依赖里的 python-pptx。"""
    return PythonPptxLoader(path)


class PythonDocxLoader:
    """轻量 DOCX loader，读取正文段落和表格文本。"""

    def __init__(self, path: Path) -> None:
        self.path = path

    def load(self) -> list[Document]:
        if self.path.suffix.lower() != ".docx":
            raise RuntimeError(f"{self.path.suffix.lower()} 需要额外转换器，请先转为 .docx 后入库。")
        docx = DocxDocument(str(self.path))
        lines = [paragraph.text.strip() for paragraph in docx.paragraphs if paragraph.text.strip()]
        for table in docx.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        content = "\n".join(lines).strip()
        return [Document(page_content=content, metadata={"file_type": ".docx"})] if content else []


class PyMuPdfLoader:
    """轻量 PDF loader，按页提取文本层内容。"""

    def __init__(self, path: Path) -> None:
        self.path = path

    def load(self) -> list[Document]:
        documents: list[Document] = []
        with fitz.open(str(self.path)) as pdf:
            for page_index, page in enumerate(pdf):
                text = page.get_text("text").strip()
                if text:
                    documents.append(Document(page_content=text, metadata={"page": page_index, "file_type": ".pdf"}))
        return documents


class PythonPptxLoader:
    """轻量 PPTX loader，按幻灯片提取文本框和表格文本。"""

    def __init__(self, path: Path) -> None:
        self.path = path

    def load(self) -> list[Document]:
        if self.path.suffix.lower() != ".pptx":
            raise RuntimeError(f"{self.path.suffix.lower()} 需要额外转换器，请先转为 .pptx 后入库。")
        presentation = Presentation(str(self.path))
        documents: list[Document] = []
        for slide_index, slide in enumerate(presentation.slides):
            lines: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            lines.append(" | ".join(cells))
                text = getattr(shape, "text", "").strip()
                if text:
                    lines.append(text)
            content = "\n".join(dict.fromkeys(lines)).strip()
            if content:
                documents.append(Document(page_content=content, metadata={"page": slide_index, "file_type": ".pptx"}))
        return documents

class TableDocumentLoader:
    """把 CSV/Excel 表格包装成 LangChain loader 协议。表格行转换成 Document 后复用主链路。"""

    def __init__(self, path: Path) -> None:
        """保存待加载的表格路径。"""
        self.path = path

    def load(self) -> list[Document]:
        """加载表格并返回行级 Document。"""
        return load_table_file(self.path)

def _table_loader(path: Path) -> SupportsLoad:
    """创建表格资料加载器。"""
    return TableDocumentLoader(path)


DOCUMENT_LOADER_SPECS: tuple[DocumentLoaderSpec, ...] = (
    DocumentLoaderSpec(
        suffixes=(".txt", ".md"),
        factory=_utf8_text_loader,
        description="UTF-8 文本和 Markdown；Markdown 保留原文给标题切分器处理。",
    ),
    DocumentLoaderSpec(
        suffixes=(".pdf",),
        factory=_pdf_loader,
        description="PDF 文本层解析；扫描件 OCR 不默认进入主链路。",
    ),
    DocumentLoaderSpec(
        suffixes=(".docx", ".doc"),
        factory=_word_loader,
        description="Word 文档文本解析。",
    ),
    DocumentLoaderSpec(
        suffixes=(".ppt", ".pptx"),
        factory=_powerpoint_loader,
        description="PowerPoint 文本解析。",
    ),
    DocumentLoaderSpec(
        suffixes=(".csv", ".xlsx", ".xls"),
        factory=_table_loader,
        description="CSV/Excel 表格解析；按行保留表头、sheet 和单元格键值。",
    ),
)

DOCUMENT_LOADER_REGISTRY: dict[str, DocumentLoaderSpec] = {
    suffix: spec
    for spec in DOCUMENT_LOADER_SPECS
    for suffix in spec.suffixes
}
SUPPORTED_DOCUMENT_SUFFIXES = tuple(sorted(DOCUMENT_LOADER_REGISTRY))

def get_document_loader_spec(path: Path) -> DocumentLoaderSpec | None:
    """根据文件后缀获取加载器注册项。后缀标准化集中处理，避免散落。"""
    return DOCUMENT_LOADER_REGISTRY.get(path.suffix.lower())

def load_file(path: Path) -> list[Document]:
    """把一个受支持的本地文件加载为 LangChain Document 对象。返回值统一为 Document，后续不关心原文件格式。"""
    spec = get_document_loader_spec(path)
    if spec is None:
        raise ValueError(f"不支持的文档类型：{path}")
    logger.debug("Loading document with %s: %s", spec.description, path)
    loader = spec.create_loader(path)
    return loader.load()


